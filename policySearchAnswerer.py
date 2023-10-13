import collections
import collections.abc
import os
import fnmatch
from pptx import Presentation
from transformers import pipeline, AutoTokenizer, AutoModelForQuestionAnswering
import torch


def read_powerpoint(file_path):

    presentation = Presentation(file_path)

    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + " "

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\t"
                text += "\n"

    return text


# Identifying Powerpoints in the Paradox 2023 Policy folder
def powerpoint_finder():
    matched_files = []
    for root, dirs, files in os.walk("/Users/will.richman/Paradox2023Policies"):
        for filename in fnmatch.filter(files, "*.pptx"):
            matched_files.append(os.path.join(root, filename))

    return matched_files


# Leveraging Q&A model to work through our policies
def policy_search(question):

    #
    file_list = powerpoint_finder()
    complete_policies = ""

    model_checkpoint = "consciousAI/question-answering-roberta-base-s-v2"

    # Read the powerpoints and gather all the information as context for the model
    for file in file_list:
        text = read_powerpoint(file)
        complete_policies += text

    context = complete_policies

    question_answerer = pipeline("question-answering", model=model_checkpoint)
    answer = question_answerer(question=question, context=context)

    return question, answer['answer'], answer['score']